"""
PPT Analyzer for AutoDoc AI

This module analyzes PowerPoint presentations to automatically detect:
- Model type (frequency, severity, loss_ratio, retention, etc.)
- Model year
- Model purpose/description

Usage:
    from utils.ppt_analyzer import analyze_ppt
    
    result = analyze_ppt("path/to/model.pptx")
    print(f"Detected: {result['model_type']} model from {result['year']}")
"""

from pathlib import Path
from typing import Dict, Optional, List, Tuple
import re
from datetime import datetime
import logging

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not installed - PPT analysis disabled")

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


class PPTAnalysisResult:
    """Result of PPT analysis."""
    
    def __init__(
        self,
        model_type: str,
        model_type_confidence: float,
        year: Optional[int],
        year_confidence: float,
        detected_keywords: Dict[str, int],
        slide_count: int,
        text_length: int,
        title: Optional[str] = None
    ):
        self.model_type = model_type
        self.model_type_confidence = model_type_confidence
        self.year = year
        self.year_confidence = year_confidence
        self.detected_keywords = detected_keywords
        self.slide_count = slide_count
        self.text_length = text_length
        self.title = title
    
    def __repr__(self):
        return (
            f"PPTAnalysisResult(model_type={self.model_type}, "
            f"confidence={self.model_type_confidence:.2f}, "
            f"year={self.year})"
        )
    
    def to_dict(self) -> Dict:
        """Convert to dictionary."""
        return {
            'model_type': self.model_type,
            'model_type_confidence': self.model_type_confidence,
            'year': self.year,
            'year_confidence': self.year_confidence,
            'detected_keywords': self.detected_keywords,
            'slide_count': self.slide_count,
            'text_length': self.text_length,
            'title': self.title
        }


# Model type detection patterns
MODEL_TYPE_PATTERNS = {
    'frequency': [
        r'\bfrequency\b',
        r'\bclaim\s+count\b',
        r'\bclaim\s+occurrence\b',
        r'\bpolicy\s+count\b',
        r'\blogistic\s+regression\b',
        r'\bbinary\s+classification\b',
        r'\bauc\b',
        r'\broc\s+curve\b',
        r'\bpredict\s+whether\b',
        r'\bpredict\s+if\b',
        r'\boccurrence\s+model\b'
    ],
    'severity': [
        r'\bseverity\b',
        r'\bclaim\s+amount\b',
        r'\bclaim\s+cost\b',
        r'\bclaim\s+size\b',
        r'\baverage\s+claim\b',
        r'\bgamma\s+regression\b',
        r'\bgamma\s+distribution\b',
        r'\bglm\s+gamma\b',
        r'\blog\s+link\b',
        r'\bmean\s+absolute\s+error\b',
        r'\bmae\b',
        r'\brmse\b'
    ],
    'pure_premium': [
        r'\bpure\s+premium\b',
        r'\btotal\s+loss\s+cost\b',
        r'\bfrequency\s*\*\s*severity\b',
        r'\bfrequency\s+times\s+severity\b',
        r'\btweedie\b'
    ],
    'loss_ratio': [
        r'\bloss\s+ratio\b',
        r'\bcombined\s+ratio\b',
        r'\bincurred\s+losses\b',
        r'\bearned\s+premium\b',
        r'\bultimate\s+loss\s+ratio\b'
    ],
    'retention': [
        r'\bretention\b',
        r'\bchurn\b',
        r'\battrition\b',
        r'\blapse\b',
        r'\brenewal\b',
        r'\bsurvival\s+analysis\b',
        r'\bcox\s+proportional\b'
    ],
    'pricing': [
        r'\bpricing\s+model\b',
        r'\brate\s+indication\b',
        r'\bbase\s+rate\b',
        r'\brating\s+algorithm\b',
        r'\bpremium\s+calculation\b'
    ]
}


def extract_all_text(ppt_path: str) -> Tuple[str, Optional[str], int]:
    """
    Extract all text from PowerPoint presentation.
    
    Args:
        ppt_path: Path to PPT file
        
    Returns:
        Tuple of (all_text, title, slide_count)
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")
    
    try:
        prs = Presentation(ppt_path)
        all_text = []
        title = None
        
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    all_text.append(text)
                    
                    # First slide, first text shape is usually the title
                    if i == 0 and title is None and len(text) < 100:
                        title = text
        
        combined_text = " ".join(all_text).lower()
        slide_count = len(prs.slides)
        
        logger.info(f"Extracted {len(combined_text)} chars from {slide_count} slides")
        return combined_text, title, slide_count
        
    except Exception as e:
        logger.error(f"Error extracting text from PPT: {e}")
        raise


def detect_model_type(text: str) -> Tuple[str, float, Dict[str, int]]:
    """
    Detect model type from text content.
    
    Args:
        text: Full text content (lowercase)
        
    Returns:
        Tuple of (model_type, confidence, keyword_counts)
    """
    keyword_counts = {}
    
    # Count matches for each model type
    for model_type, patterns in MODEL_TYPE_PATTERNS.items():
        count = 0
        for pattern in patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            count += len(matches)
        keyword_counts[model_type] = count
    
    # Find model type with most matches
    if not keyword_counts or max(keyword_counts.values()) == 0:
        logger.warning("No model type keywords found in PPT")
        return "unknown", 0.0, keyword_counts
    
    best_model = max(keyword_counts.items(), key=lambda x: x[1])
    model_type = best_model[0]
    count = best_model[1]
    
    # Calculate confidence based on:
    # - Number of matches
    # - Distinctiveness (how much better than 2nd place)
    sorted_counts = sorted(keyword_counts.values(), reverse=True)
    second_best = sorted_counts[1] if len(sorted_counts) > 1 else 0
    
    # Confidence factors
    match_score = min(count / 10.0, 1.0)  # 10+ matches = full confidence
    distinctiveness = (count - second_best) / max(count, 1)
    
    confidence = (match_score * 0.7 + distinctiveness * 0.3)
    
    logger.info(f"Detected model type: {model_type} (confidence: {confidence:.2f}, matches: {count})")
    return model_type, confidence, keyword_counts


def detect_year(text: str) -> Tuple[Optional[int], float]:
    """
    Detect model year from text content.
    
    Args:
        text: Full text content
        
    Returns:
        Tuple of (year, confidence)
    """
    current_year = datetime.now().year
    
    # Look for 4-digit years near keywords like "model", "update", "version"
    year_patterns = [
        r'(?:model|update|version|year|q[1-4])\s+(\d{4})',
        r'(\d{4})\s+(?:model|update|version)',
        r'copyright\s+(\d{4})',
        r'©\s*(\d{4})'
    ]
    
    found_years = []
    for pattern in year_patterns:
        matches = re.findall(pattern, text, re.IGNORECASE)
        for match in matches:
            year = int(match)
            # Reasonable range: 2015-2030
            if 2015 <= year <= 2030:
                found_years.append(year)
    
    if not found_years:
        # Try finding any 4-digit number that looks like a year
        all_years = re.findall(r'\b(20[1-3][0-9])\b', text)
        found_years = [int(y) for y in all_years if 2015 <= int(y) <= 2030]
    
    if not found_years:
        logger.warning("No year detected in PPT")
        return None, 0.0
    
    # Most common year wins
    from collections import Counter
    year_counts = Counter(found_years)
    most_common_year, count = year_counts.most_common(1)[0]
    
    # Confidence based on:
    # - How recent (prefer recent years)
    # - How often mentioned
    recency_score = 1.0 - (abs(current_year - most_common_year) / 10.0)
    recency_score = max(0.0, min(1.0, recency_score))
    
    mention_score = min(count / 5.0, 1.0)  # 5+ mentions = full confidence
    
    confidence = (recency_score * 0.6 + mention_score * 0.4)
    
    logger.info(f"Detected year: {most_common_year} (confidence: {confidence:.2f}, mentions: {count})")
    return most_common_year, confidence


def analyze_ppt(ppt_path: str) -> PPTAnalysisResult:
    """
    Analyze PowerPoint presentation to detect model type and year.
    
    Args:
        ppt_path: Path to PPT file
        
    Returns:
        PPTAnalysisResult with detected information
        
    Example:
        >>> result = analyze_ppt("collision_severity_model.pptx")
        >>> print(f"Model: {result.model_type} ({result.model_type_confidence:.0%})")
        Model: severity (85%)
        >>> print(f"Year: {result.year} ({result.year_confidence:.0%})")
        Year: 2024 (92%)
    """
    logger.info(f"Analyzing PPT: {ppt_path}")
    
    # Extract text
    text, title, slide_count = extract_all_text(ppt_path)
    text_length = len(text)
    
    # Detect model type
    model_type, model_confidence, keyword_counts = detect_model_type(text)
    
    # Detect year
    year, year_confidence = detect_year(text)
    
    result = PPTAnalysisResult(
        model_type=model_type,
        model_type_confidence=model_confidence,
        year=year,
        year_confidence=year_confidence,
        detected_keywords=keyword_counts,
        slide_count=slide_count,
        text_length=text_length,
        title=title
    )
    
    logger.info(f"Analysis complete: {result}")
    return result


def format_confidence(confidence: float) -> str:
    """Format confidence score as human-readable string."""
    if confidence >= 0.8:
        return "HIGH"
    elif confidence >= 0.5:
        return "MEDIUM"
    else:
        return "LOW"


if __name__ == "__main__":
    # Example usage
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python ppt_analyzer.py <path_to_ppt>")
        sys.exit(1)
    
    ppt_path = sys.argv[1]
    result = analyze_ppt(ppt_path)
    
    print("\n" + "=" * 60)
    print("PPT ANALYSIS RESULTS")
    print("=" * 60)
    print(f"File: {ppt_path}")
    print(f"Slides: {result.slide_count}")
    print(f"Text Length: {result.text_length:,} chars")
    if result.title:
        print(f"Title: {result.title}")
    print(f"\nModel Type: {result.model_type}")
    print(f"Confidence: {format_confidence(result.model_type_confidence)} ({result.model_type_confidence:.0%})")
    print(f"\nYear: {result.year}")
    print(f"Confidence: {format_confidence(result.year_confidence)} ({result.year_confidence:.0%})")
    print(f"\nKeyword Matches:")
    for model_type, count in sorted(result.detected_keywords.items(), key=lambda x: x[1], reverse=True):
        if count > 0:
            print(f"  {model_type}: {count}")
    print("=" * 60)
