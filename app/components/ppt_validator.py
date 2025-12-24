"""
PowerPoint Validator Component for AutoDoc AI

Validates uploaded PowerPoint files for:
- File format and structure
- Required slides and content
- Model presentation template compliance
"""

import streamlit as st
from pathlib import Path
from typing import Dict, List, Tuple, Optional
from dataclasses import dataclass, field
from pptx import Presentation
from pptx.exc import PackageNotFoundError
import io


@dataclass
class ValidationResult:
    """Result of PPT validation."""
    is_valid: bool
    errors: List[str] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)
    info: List[str] = field(default_factory=list)
    slide_count: int = 0
    detected_sections: List[str] = field(default_factory=list)


class PPTValidator:
    """
    Validator for insurance model presentation PowerPoints.

    Checks for:
    - File integrity
    - Minimum slide requirements
    - Recommended section structure
    - Model-specific content
    """

    # Expected sections in a model presentation
    RECOMMENDED_SECTIONS = [
        "Executive Summary",
        "Business Context",
        "Data Environment",
        "Methodology",
        "Model Results",
        "Validation",
        "Implementation",
        "Monitoring"
    ]

    # Minimum requirements
    MIN_SLIDES = 10
    MAX_SLIDES = 30

    def __init__(self):
        """Initialize PPT validator."""
        pass

    def validate_file(self, file_bytes: bytes, filename: str) -> ValidationResult:
        """
        Validate a PowerPoint file.

        Args:
            file_bytes: Bytes of the uploaded file
            filename: Original filename

        Returns:
            ValidationResult with detailed findings
        """
        result = ValidationResult(is_valid=True)

        # Check file extension
        if not filename.lower().endswith(('.pptx', '.ppt')):
            result.errors.append(
                f"Invalid file format: {filename}. Expected .pptx or .ppt"
            )
            result.is_valid = False
            return result

        # Try to load the presentation
        try:
            file_stream = io.BytesIO(file_bytes)
            prs = Presentation(file_stream)
        except PackageNotFoundError:
            result.errors.append(
                "Cannot read PowerPoint file. File may be corrupted or in an unsupported format."
            )
            result.is_valid = False
            return result
        except Exception as e:
            result.errors.append(f"Error reading PowerPoint: {str(e)}")
            result.is_valid = False
            return result

        # Validate slide count
        slide_count = len(prs.slides)
        result.slide_count = slide_count

        if slide_count < self.MIN_SLIDES:
            result.errors.append(
                f"Too few slides: {slide_count} found, minimum {self.MIN_SLIDES} expected"
            )
            result.is_valid = False
        elif slide_count > self.MAX_SLIDES:
            result.warnings.append(
                f"Many slides: {slide_count} found. Recommended maximum is {self.MAX_SLIDES}"
            )
        else:
            result.info.append(f"Slide count: {slide_count} (✓ within recommended range)")

        # Detect sections from slide titles
        detected_sections = self._extract_slide_titles(prs)
        result.detected_sections = detected_sections

        # Check for recommended sections
        missing_sections = self._check_sections(detected_sections)
        if missing_sections:
            result.warnings.append(
                f"Missing recommended sections: {', '.join(missing_sections)}"
            )
        else:
            result.info.append("✓ All recommended sections detected")

        # Check for model-specific keywords
        model_keywords = self._check_model_keywords(prs)
        if model_keywords:
            result.info.append(
                f"Detected model keywords: {', '.join(model_keywords[:5])}"
            )
        else:
            result.warnings.append(
                "No model-specific keywords detected (GLM, frequency, severity, etc.)"
            )

        # Check for tables (important for model results)
        table_count = self._count_tables(prs)
        if table_count > 0:
            result.info.append(f"✓ Found {table_count} tables (likely model results)")
        else:
            result.warnings.append("No tables detected. Model results typically include tables.")

        return result

    def _extract_slide_titles(self, prs: Presentation) -> List[str]:
        """Extract titles from all slides."""
        titles = []
        for slide in prs.slides:
            if slide.shapes.title and slide.shapes.title.text:
                titles.append(slide.shapes.title.text.strip())
        return titles

    def _check_sections(self, detected_titles: List[str]) -> List[str]:
        """
        Check which recommended sections are missing.

        Returns list of missing section names.
        """
        # Convert to lowercase for fuzzy matching
        detected_lower = [t.lower() for t in detected_titles]

        missing = []
        for section in self.RECOMMENDED_SECTIONS:
            # Check if any slide title contains the section name (fuzzy match)
            found = any(section.lower() in title for title in detected_lower)
            if not found:
                missing.append(section)

        return missing

    def _check_model_keywords(self, prs: Presentation) -> List[str]:
        """
        Check for model-specific keywords in the presentation.

        Returns list of detected keywords.
        """
        keywords = {
            "glm", "frequency", "severity", "model", "regression",
            "xgboost", "gradient boosting", "neural network",
            "auc", "gini", "r-squared", "rmse", "validation",
            "territory", "rating", "premium", "loss ratio",
            "claim", "exposure", "coverage", "risk"
        }

        detected = set()

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_lower = shape.text.lower()
                    for keyword in keywords:
                        if keyword in text_lower:
                            detected.add(keyword)

        return sorted(list(detected))

    def _count_tables(self, prs: Presentation) -> int:
        """Count tables in the presentation."""
        table_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 19:  # Table shape type
                    table_count += 1
        return table_count

    def display_validation_results(self, result: ValidationResult):
        """
        Display validation results in Streamlit.

        Args:
            result: ValidationResult to display
        """
        if result.is_valid:
            st.success("✅ PowerPoint validation passed!")
        else:
            st.error("❌ PowerPoint validation failed")

        # Show errors (if any)
        if result.errors:
            st.subheader("🔴 Errors")
            for error in result.errors:
                st.error(error)

        # Show warnings (if any)
        if result.warnings:
            st.subheader("⚠️ Warnings")
            for warning in result.warnings:
                st.warning(warning)

        # Show info messages
        if result.info:
            st.subheader("ℹ️ Information")
            for info in result.info:
                st.info(info)

        # Show detected sections
        if result.detected_sections:
            st.subheader("📄 Detected Slides")
            with st.expander("Show all slide titles"):
                for i, title in enumerate(result.detected_sections, 1):
                    st.write(f"{i}. {title}")


def validate_and_display(file_bytes: bytes, filename: str) -> Tuple[bool, ValidationResult]:
    """
    Validate a PowerPoint file and display results.

    Args:
        file_bytes: File bytes from uploaded file
        filename: Original filename

    Returns:
        Tuple of (is_valid, ValidationResult)
    """
    validator = PPTValidator()
    result = validator.validate_file(file_bytes, filename)
    validator.display_validation_results(result)

    return result.is_valid, result


def show_template_info():
    """Display information about the recommended PowerPoint template."""
    st.info(
        """
        📋 **Recommended Presentation Structure**

        Your model presentation should include these sections:

        1. **Executive Summary** - Key findings and recommendations
        2. **Business Context** - Problem statement and objectives
        3. **Data Environment** - Data sources, quality, variables
        4. **Methodology** - Model type, technique, assumptions
        5. **Model Results** - Performance metrics, coefficients
        6. **Validation** - Out-of-sample testing, lift charts
        7. **Implementation** - Deployment plan and timeline
        8. **Monitoring** - Ongoing tracking and governance

        **Requirements:**
        - 10-30 slides (recommended 15-20)
        - Include tables with model results
        - Clear slide titles for each section
        """
    )


if __name__ == "__main__":
    # Test the PPT validator
    st.set_page_config(page_title="PPT Validator Test", layout="wide")

    st.title("PowerPoint Validator Test")

    show_template_info()

    st.markdown("---")

    # File uploader for testing
    uploaded_file = st.file_uploader(
        "Upload a PowerPoint presentation",
        type=['pptx', 'ppt'],
        help="Upload a model presentation to validate"
    )

    if uploaded_file is not None:
        st.subheader("Validation Results")

        file_bytes = uploaded_file.read()
        is_valid, result = validate_and_display(file_bytes, uploaded_file.name)

        st.markdown("---")
        st.subheader("Raw Validation Result")
        st.json({
            "is_valid": result.is_valid,
            "slide_count": result.slide_count,
            "errors_count": len(result.errors),
            "warnings_count": len(result.warnings),
            "info_count": len(result.info),
            "detected_sections_count": len(result.detected_sections)
        })
