"""
PowerPoint Parser for AutoDoc AI

Extracts structured content from PowerPoint presentations:
- Slide titles and content
- Text blocks
- Tables
- Images (metadata)
- Structure and organization
"""

from pathlib import Path
from typing import Dict, List, Optional
from dataclasses import dataclass, field
import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class SlideContent:
    """Content extracted from a single slide."""
    slide_number: int
    title: str
    text_content: List[str] = field(default_factory=list)
    tables: List[Dict] = field(default_factory=list)
    has_chart: bool = False
    has_image: bool = False
    notes: str = ""


@dataclass
class PPTContent:
    """Complete content extracted from PowerPoint."""
    filename: str
    total_slides: int
    slides: List[SlideContent] = field(default_factory=list)
    total_text_blocks: int = 0
    total_tables: int = 0
    total_charts: int = 0
    total_images: int = 0


class PPTXParser:
    """
    Parser for PowerPoint (.pptx) files.

    Extracts structured content for use in documentation generation.
    """

    def __init__(self):
        """Initialize PPTX parser."""
        pass

    def extract_from_file(self, filepath: Path) -> PPTContent:
        """
        Extract content from a PowerPoint file.

        Args:
            filepath: Path to .pptx file

        Returns:
            PPTContent with extracted data
        """
        prs = Presentation(str(filepath))
        return self._extract_content(prs, filepath.name)

    def extract_from_stream(self, file_stream: io.BytesIO) -> PPTContent:
        """
        Extract content from a file stream.

        Args:
            file_stream: BytesIO stream of .pptx file

        Returns:
            PPTContent with extracted data
        """
        prs = Presentation(file_stream)
        return self._extract_content(prs, "uploaded_file.pptx")

    def _extract_content(self, prs: Presentation, filename: str) -> PPTContent:
        """
        Extract content from a Presentation object.

        Args:
            prs: python-pptx Presentation object
            filename: Original filename

        Returns:
            PPTContent with extracted data
        """
        content = PPTContent(
            filename=filename,
            total_slides=len(prs.slides)
        )

        for i, slide in enumerate(prs.slides, 1):
            slide_content = self._extract_slide_content(slide, i)
            content.slides.append(slide_content)

            # Update totals
            content.total_text_blocks += len(slide_content.text_content)
            content.total_tables += len(slide_content.tables)
            if slide_content.has_chart:
                content.total_charts += 1
            if slide_content.has_image:
                content.total_images += 1

        return content

    def _extract_slide_content(self, slide, slide_number: int) -> SlideContent:
        """
        Extract content from a single slide.

        Args:
            slide: python-pptx Slide object
            slide_number: Slide number (1-indexed)

        Returns:
            SlideContent with extracted data
        """
        # Extract title
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        slide_content = SlideContent(
            slide_number=slide_number,
            title=title or f"Slide {slide_number}"
        )

        # Extract notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                slide_content.notes = notes_slide.notes_text_frame.text.strip()

        # Process all shapes
        for shape in slide.shapes:
            # Text content
            if hasattr(shape, "text") and shape.text.strip():
                # Skip title (already extracted)
                if shape == slide.shapes.title:
                    continue

                text = shape.text.strip()
                if text:
                    slide_content.text_content.append(text)

            # Tables
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_data = self._extract_table(shape.table)
                slide_content.tables.append(table_data)

            # Charts
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                slide_content.has_chart = True

            # Images/Pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_content.has_image = True

        return slide_content

    def _extract_table(self, table) -> Dict:
        """
        Extract data from a table shape.

        Args:
            table: python-pptx Table object

        Returns:
            Dictionary with table data
        """
        rows = []

        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            rows.append(row_data)

        return {
            "rows": len(table.rows),
            "columns": len(table.columns),
            "data": rows
        }

    def format_as_markdown(self, content: PPTContent) -> str:
        """
        Format extracted content as Markdown.

        Args:
            content: PPTContent to format

        Returns:
            Markdown-formatted string
        """
        md_lines = []

        md_lines.append(f"# {content.filename}")
        md_lines.append("")
        md_lines.append(f"**Total Slides:** {content.total_slides}")
        md_lines.append("")
        md_lines.append("---")
        md_lines.append("")

        for slide in content.slides:
            md_lines.append(f"## Slide {slide.slide_number}: {slide.title}")
            md_lines.append("")

            # Text content
            if slide.text_content:
                for text in slide.text_content:
                    # Split into bullet points if multiple paragraphs
                    paragraphs = text.split('\n')
                    for para in paragraphs:
                        if para.strip():
                            md_lines.append(f"- {para.strip()}")
                md_lines.append("")

            # Tables
            if slide.tables:
                md_lines.append("**Tables:**")
                for i, table in enumerate(slide.tables, 1):
                    md_lines.append(f"\nTable {i} ({table['rows']} rows × {table['columns']} columns)")
                    md_lines.append("")

                    # Format as Markdown table (first 2 rows only for preview)
                    if table['data']:
                        for row_idx, row in enumerate(table['data'][:2]):
                            md_lines.append("| " + " | ".join(row) + " |")
                            if row_idx == 0:
                                md_lines.append("| " + " | ".join(["---"] * len(row)) + " |")

                        if len(table['data']) > 2:
                            md_lines.append(f"\n*...{len(table['data']) - 2} more rows*")
                md_lines.append("")

            # Charts/Images
            if slide.has_chart:
                md_lines.append("📊 *Contains chart(s)*")
                md_lines.append("")

            if slide.has_image:
                md_lines.append("🖼️ *Contains image(s)*")
                md_lines.append("")

            # Notes
            if slide.notes:
                md_lines.append("**Speaker Notes:**")
                md_lines.append(f"> {slide.notes}")
                md_lines.append("")

            md_lines.append("---")
            md_lines.append("")

        return "\n".join(md_lines)

    def get_summary_stats(self, content: PPTContent) -> Dict:
        """
        Get summary statistics from extracted content.

        Args:
            content: PPTContent to analyze

        Returns:
            Dictionary with summary statistics
        """
        total_words = 0
        total_bullets = 0

        for slide in content.slides:
            for text in slide.text_content:
                total_words += len(text.split())
                total_bullets += text.count('\n') + 1

        return {
            "filename": content.filename,
            "total_slides": content.total_slides,
            "total_text_blocks": content.total_text_blocks,
            "total_words": total_words,
            "total_bullets": total_bullets,
            "total_tables": content.total_tables,
            "total_charts": content.total_charts,
            "total_images": content.total_images,
            "avg_words_per_slide": total_words / content.total_slides if content.total_slides > 0 else 0
        }


def main():
    """Test the PPTX parser."""
    import sys

    if len(sys.argv) < 2:
        print("Usage: python pptx_parser.py <path_to_pptx>")
        print("\nTest mode: Creating sample output")

        # Demo mode - show what it would extract
        print("\n=== PPTX Parser Demo ===\n")
        print("This parser extracts:")
        print("1. Slide titles and content")
        print("2. Text blocks and bullet points")
        print("3. Tables (with full data)")
        print("4. Chart and image presence")
        print("5. Speaker notes")
        print("\nSample output structure:")
        print("""
        PPTContent(
            filename='model_presentation.pptx',
            total_slides=15,
            slides=[
                SlideContent(
                    slide_number=1,
                    title='Executive Summary',
                    text_content=['Key findings...', 'Recommendations...'],
                    tables=[],
                    has_chart=False,
                    has_image=True
                ),
                ...
            ]
        )
        """)
        return

    filepath = Path(sys.argv[1])

    if not filepath.exists():
        print(f"Error: File not found: {filepath}")
        return

    print(f"Parsing PowerPoint: {filepath}")
    print("=" * 60)

    parser = PPTXParser()
    content = parser.extract_from_file(filepath)

    # Print summary
    stats = parser.get_summary_stats(content)
    print("\n=== Summary Statistics ===")
    for key, value in stats.items():
        print(f"{key}: {value}")

    # Print Markdown output
    print("\n=== Markdown Output ===\n")
    markdown = parser.format_as_markdown(content)
    print(markdown)


if __name__ == "__main__":
    main()
