"""
Create Example PowerPoint Presentations for AutoDoc AI

Generates 3 complete actuarial model presentations:
1. Bodily Injury Frequency Model (GLM Poisson)
2. Collision Severity Model (GLM Gamma + Telematics)
3. Comprehensive Coverage Model (XGBoost)

Each presentation is 15-20 slides with tables, charts, and comprehensive content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
from datetime import datetime


def create_title_slide(prs, title, subtitle):
    """Create a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def create_content_slide(prs, title):
    """Create a slide with title and content area."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    return slide


def create_blank_slide(prs):
    """Create a blank slide."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_table(slide, rows, cols, left, top, width, height, data):
    """Add a table to a slide with data."""
    # Use actual data length if it differs from specified rows
    actual_rows = len(data)
    if actual_rows != rows:
        rows = actual_rows

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Fill table with data first
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.rows[row_idx].cells[col_idx]
            cell.text = str(cell_data)

    # Style header row
    for col_idx in range(cols):
        cell = table.rows[0].cells[col_idx]
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(31, 119, 180)  # Blue header
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
        paragraph.font.bold = True
        paragraph.font.size = Pt(11)

    # Alternate row colors for data rows
    for row_idx in range(1, rows):
        if row_idx % 2 == 0:
            for col_idx in range(cols):
                cell = table.rows[row_idx].cells[col_idx]
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

    return table


def add_bullet_points(text_frame, points):
    """Add bullet points to a text frame."""
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(14)


def create_frequency_model_ppt():
    """
    Create Bodily Injury Frequency Model presentation (GLM Poisson).

    Based on 2022_frequency_model_doc.md:
    - GLM Poisson regression
    - 12 predictor variables
    - AUC 0.72, Gini 0.44
    """
    print("Creating Bodily Injury Frequency Model PPT...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = create_title_slide(
        prs,
        "Bodily Injury Frequency Model",
        "GLM Poisson Regression Model\n2024 Model Development\nPresented to Model Risk Committee"
    )

    # Slide 2: Executive Summary
    slide = create_content_slide(prs, "Executive Summary")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Developed GLM Poisson model for Bodily Injury claim frequency prediction",
        "12 statistically significant predictors selected (p < 0.05)",
        "Model Performance: AUC 0.72, Gini 0.44 on holdout test set",
        "15% lift in top decile vs. random selection",
        "Passes all regulatory requirements (NAIC Model Audit Rule, ASOPs)",
        "Recommendation: Approve for production deployment Q3 2024"
    ])

    # Slide 3: Business Context
    slide = create_content_slide(prs, "Business Context & Objectives")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Current model developed in 2019, performance degrading over time",
        "Need to incorporate new risk factors (telematics, credit-based insurance scores)",
        "Regulatory requirement for annual model validation and refresh",
        "Objective: Improve risk segmentation by 10-15% vs. current champion",
        "Support rate adequacy and competitiveness in key markets",
        "Ensure compliance with evolving regulatory standards"
    ])

    # Slide 4: Data Overview
    slide = create_content_slide(prs, "Data Overview")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Data Source: PolicyMaster system (2019-2023 experience period)",
        "Sample Size: 500,000 policies, 50,000 claims",
        "Geography: 15 states (Midwest and Southeast regions)",
        "Coverage: Personal Auto Bodily Injury Liability",
        "Data Quality: 98.5% complete, extensive validation performed"
    ])

    # Add data summary table
    data = [
        ["Metric", "Training", "Validation", "Test"],
        ["Policies", "300,000", "100,000", "100,000"],
        ["Claims", "30,000", "10,000", "10,000"],
        ["Claim Frequency", "10.0%", "10.1%", "9.9%"],
        ["Time Period", "2019-2021", "2022", "2023"],
        ["Data Quality", "98.7%", "98.5%", "98.3%"]
    ]
    add_table(slide, 6, 4, Inches(0.5), Inches(3.5), Inches(9), Inches(2.5), data)

    # Slide 5: Descriptive Statistics
    slide = create_content_slide(prs, "Descriptive Statistics - Continuous Variables")
    data = [
        ["Variable", "Mean", "Std Dev", "Min", "Max", "Missing %"],
        ["Driver Age", "42.3", "14.2", "16", "95", "0.1%"],
        ["Vehicle Age", "7.8", "4.5", "0", "25", "0.3%"],
        ["Annual Miles", "12,450", "5,200", "1,000", "50,000", "2.1%"],
        ["Years Licensed", "20.5", "13.8", "0", "70", "0.5%"],
        ["Credit Score", "685", "95", "400", "850", "4.2%"]
    ]
    add_table(slide, 6, 6, Inches(0.5), Inches(1.8), Inches(9), Inches(3), data)

    # Slide 6: Descriptive Statistics - Categorical
    slide = create_content_slide(prs, "Descriptive Statistics - Categorical Variables")
    data = [
        ["Variable", "Level", "Count", "Percentage", "Avg Frequency"],
        ["Territory", "Urban High", "75,000", "15%", "14.2%"],
        ["Territory", "Urban Medium", "100,000", "20%", "11.5%"],
        ["Territory", "Suburban", "225,000", "45%", "9.2%"],
        ["Territory", "Rural", "100,000", "20%", "7.8%"],
        ["Vehicle Type", "Sedan", "250,000", "50%", "9.5%"],
        ["Vehicle Type", "SUV", "150,000", "30%", "10.8%"],
        ["Vehicle Type", "Truck", "100,000", "20%", "11.2%"]
    ]
    add_table(slide, 8, 5, Inches(0.5), Inches(1.8), Inches(9), Inches(4), data)

    # Slide 7: Variable Selection
    slide = create_content_slide(prs, "Variable Selection Process")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Started with 45 candidate variables from business knowledge",
        "Applied correlation analysis to remove multicollinearity (VIF < 5)",
        "Used stepwise selection (AIC criterion) for initial model",
        "Validated business intuition and actuarial judgment",
        "Removed variables with p-value > 0.05",
        "Final model: 12 statistically significant predictors"
    ])

    # Slide 8: Selected Variables
    slide = create_content_slide(prs, "Final Variable Set with Statistical Significance")
    data = [
        ["Variable", "Type", "Levels/Range", "Coefficient", "Std Error", "p-value"],
        ["Intercept", "-", "-", "-2.450", "0.025", "< 0.001"],
        ["Driver Age", "Continuous", "16-95", "-0.012", "0.002", "< 0.001"],
        ["Vehicle Age", "Continuous", "0-25", "0.025", "0.004", "< 0.001"],
        ["Territory", "Categorical", "15 zones", "varies", "-", "< 0.001"],
        ["Prior Claims (1)", "Binary", "0/1", "0.485", "0.035", "< 0.001"],
        ["Prior Claims (2+)", "Binary", "0/1", "0.852", "0.048", "< 0.001"],
        ["Annual Miles", "Continuous", "1K-50K", "0.00003", "0.00001", "0.002"],
        ["Vehicle Type", "Categorical", "3 types", "varies", "-", "< 0.001"],
        ["Coverage Limit", "Categorical", "5 limits", "varies", "-", "0.012"],
        ["Credit Score Band", "Categorical", "5 bands", "varies", "-", "< 0.001"],
        ["Years Licensed", "Continuous", "0-70", "-0.008", "0.003", "0.005"],
        ["Gender", "Binary", "M/F", "0.125", "0.028", "< 0.001"],
        ["Multi-Policy", "Binary", "Y/N", "-0.185", "0.032", "< 0.001"]
    ]
    add_table(slide, 13, 6, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5), data)

    # Slide 9: Single Factor Analysis - Driver Age
    slide = create_content_slide(prs, "Single Factor Analysis: Driver Age")
    data = [
        ["Age Band", "Policies", "Claims", "Frequency", "Relativity", "Credibility"],
        ["16-20", "15,000", "2,400", "16.0%", "1.60", "High"],
        ["21-25", "35,000", "4,550", "13.0%", "1.30", "High"],
        ["26-30", "45,000", "4,950", "11.0%", "1.10", "High"],
        ["31-40", "90,000", "9,000", "10.0%", "1.00", "High"],
        ["41-50", "120,000", "10,800", "9.0%", "0.90", "High"],
        ["51-60", "110,000", "8,800", "8.0%", "0.80", "High"],
        ["61-70", "60,000", "4,800", "8.0%", "0.80", "High"],
        ["71+", "25,000", "2,750", "11.0%", "1.10", "Medium"]
    ]
    add_table(slide, 9, 6, Inches(1), Inches(1.8), Inches(8), Inches(4.5), data)

    # Slide 10: Single Factor Analysis - Territory
    slide = create_content_slide(prs, "Single Factor Analysis: Territory")
    data = [
        ["Territory", "Description", "Policies", "Claims", "Frequency", "Relativity"],
        ["Zone 1", "Urban - Downtown", "25,000", "3,550", "14.2%", "1.42"],
        ["Zone 2", "Urban - High Density", "50,000", "5,750", "11.5%", "1.15"],
        ["Zone 3-5", "Suburban - High", "75,000", "7,500", "10.0%", "1.00"],
        ["Zone 6-10", "Suburban - Medium", "150,000", "13,800", "9.2%", "0.92"],
        ["Zone 11-12", "Suburban - Low", "100,000", "8,500", "8.5%", "0.85"],
        ["Zone 13-15", "Rural", "100,000", "7,800", "7.8%", "0.78"]
    ]
    add_table(slide, 7, 6, Inches(0.8), Inches(1.8), Inches(8.4), Inches(4), data)

    # Slide 11: Single Factor Analysis - Prior Claims
    slide = create_content_slide(prs, "Single Factor Analysis: Prior Claims History")
    data = [
        ["Prior Claims", "Policies", "Claims", "Frequency", "Relativity", "Statistical Sig"],
        ["0 claims", "400,000", "36,000", "9.0%", "1.00", "Base"],
        ["1 claim", "80,000", "11,920", "14.9%", "1.66", "p < 0.001"],
        ["2 claims", "15,000", "2,925", "19.5%", "2.17", "p < 0.001"],
        ["3+ claims", "5,000", "1,200", "24.0%", "2.67", "p < 0.001"]
    ]
    add_table(slide, 5, 6, Inches(1.2), Inches(2), Inches(7.6), Inches(3), data)

    tf = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(8), Inches(1)).text_frame
    p = tf.paragraphs[0]
    p.text = "Strong persistence effect: Drivers with prior claims show 66-167% higher frequency"
    p.font.size = Pt(12)
    p.font.italic = True

    # Slide 12: Model Development Approach
    slide = create_content_slide(prs, "Model Development Methodology")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Model Type: Generalized Linear Model (GLM) with Poisson distribution",
        "Link Function: Log link for multiplicative effects",
        "Software: R statistical package (glm function)",
        "Training/Validation/Test Split: 60% / 20% / 20%",
        "Cross-validation: 5-fold CV on training set for hyperparameter tuning",
        "Assumptions validated: independence, linearity in log scale, equidispersion",
        "Regularization: None (traditional GLM approach)",
        "Missing data handling: Imputation with separate category for categoricals"
    ])

    # Slide 13: Champion vs. Challenger Comparison
    slide = create_content_slide(prs, "Model Comparison: Champion vs. Challenger")
    data = [
        ["Metric", "Champion (2019)", "Challenger (2024)", "Improvement", "Status"],
        ["AUC", "0.68", "0.72", "+5.9%", "✓ Better"],
        ["Gini Coefficient", "0.36", "0.44", "+22.2%", "✓ Better"],
        ["Top Decile Lift", "2.0x", "2.3x", "+15%", "✓ Better"],
        ["Top Quintile Lift", "1.6x", "1.8x", "+12.5%", "✓ Better"],
        ["Deviance", "13,250", "12,345", "-6.8%", "✓ Better"],
        ["AIC", "148,520", "145,678", "-1.9%", "✓ Better"],
        ["# Variables", "10", "12", "+20%", "More complex"],
        ["Interpretability", "High", "High", "Equal", "✓ Good"]
    ]
    add_table(slide, 9, 5, Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5), data)

    # Slide 14: Model Results - Coefficients
    slide = create_content_slide(prs, "Model Coefficients (Selected Variables)")
    data = [
        ["Variable", "Level", "Coefficient", "Exp(Coef)", "Interpretation"],
        ["Intercept", "-", "-2.450", "0.086", "Base frequency: 8.6%"],
        ["Driver Age", "Per year", "-0.012", "0.988", "1.2% decrease per year"],
        ["Vehicle Age", "Per year", "0.025", "1.025", "2.5% increase per year"],
        ["Prior Claims", "1 claim", "0.485", "1.624", "62.4% increase"],
        ["Prior Claims", "2+ claims", "0.852", "2.344", "134.4% increase"],
        ["Territory", "Zone 1 vs 10", "0.585", "1.795", "79.5% increase"],
        ["Credit Score", "Poor vs Good", "0.425", "1.530", "53.0% increase"],
        ["Multi-Policy", "Yes vs No", "-0.185", "0.831", "16.9% decrease"]
    ]
    add_table(slide, 9, 5, Inches(0.5), Inches(1.8), Inches(9), Inches(4.5), data)

    # Slide 15: Validation Results
    slide = create_content_slide(prs, "Out-of-Sample Validation Results")

    # Validation metrics
    data = [
        ["Segment", "Policies", "Actual Freq", "Predicted Freq", "A/E Ratio", "Assessment"],
        ["Overall", "100,000", "9.9%", "10.1%", "0.98", "✓ Excellent"],
        ["Young Drivers", "15,000", "15.5%", "15.8%", "0.98", "✓ Good"],
        ["Middle Age", "60,000", "9.2%", "9.4%", "0.98", "✓ Excellent"],
        ["Senior Drivers", "25,000", "10.8%", "10.5%", "1.03", "✓ Good"],
        ["Urban Territory", "30,000", "12.5%", "12.8%", "0.98", "✓ Good"],
        ["Suburban", "50,000", "9.1%", "9.3%", "0.98", "✓ Excellent"],
        ["Rural", "20,000", "7.9%", "7.7%", "1.03", "✓ Good"]
    ]
    add_table(slide, 8, 6, Inches(0.5), Inches(1.8), Inches(9), Inches(4), data)

    tf = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.8)).text_frame
    p = tf.paragraphs[0]
    p.text = "All A/E ratios within 95-105% range. Model demonstrates excellent predictive accuracy across all segments."
    p.font.size = Pt(11)
    p.font.italic = True

    # Slide 16: Lift Chart Analysis
    slide = create_content_slide(prs, "Lift Chart: Model Discrimination")
    data = [
        ["Decile", "Policies", "Claims", "Actual Freq", "Predicted Freq", "Lift vs Random"],
        ["1 (Highest)", "10,000", "2,300", "23.0%", "23.5%", "2.32x"],
        ["2", "10,000", "1,750", "17.5%", "17.8%", "1.77x"],
        ["3", "10,000", "1,420", "14.2%", "14.5%", "1.43x"],
        ["4", "10,000", "1,180", "11.8%", "12.0%", "1.19x"],
        ["5", "10,000", "1,000", "10.0%", "10.2%", "1.01x"],
        ["6", "10,000", "850", "8.5%", "8.6%", "0.86x"],
        ["7", "10,000", "720", "7.2%", "7.3%", "0.73x"],
        ["8", "10,000", "580", "5.8%", "5.9%", "0.59x"],
        ["9", "10,000", "450", "4.5%", "4.6%", "0.45x"],
        ["10 (Lowest)", "10,000", "330", "3.3%", "3.4%", "0.33x"]
    ]
    add_table(slide, 11, 6, Inches(0.5), Inches(1.5), Inches(9), Inches(5.2), data)

    # Slide 17: Implementation Plan
    slide = create_content_slide(prs, "Implementation Plan & Timeline")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Model Approval: Model Risk Committee - June 2024",
        "System Integration: IT implementation - July 2024",
        "Parallel Testing: Shadow mode with current model - August 2024",
        "Production Deployment: Q3 2024 (September 1)",
        "Monitoring Setup: Real-time dashboards and alerts - August 2024",
        "Training: Underwriter and actuarial staff training - August 2024",
        "Documentation: Complete all regulatory documentation - July 2024",
        "First Review: 90-day post-implementation review - December 2024"
    ])

    # Slide 18: Monitoring & Governance
    slide = create_content_slide(prs, "Ongoing Monitoring & Model Governance")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Monthly Performance Monitoring:",
        "  • Actual vs. Expected frequency by segment",
        "  • Model stability metrics (PSI, CSI)",
        "  • Data quality checks",
        "Quarterly Model Review:",
        "  • Detailed validation with recent data",
        "  • Review of model assumptions",
        "  • Assessment of model performance drift",
        "Annual Model Refresh:",
        "  • Full model recalibration with latest data",
        "  • Regulatory compliance review",
        "  • Model Risk Committee approval"
    ])

    # Slide 19: Regulatory Compliance
    slide = create_content_slide(prs, "Regulatory Compliance Attestation")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "NAIC Model Audit Rule (MAR) Compliance:",
        "  • Model governance framework documented",
        "  • Independent validation completed",
        "  • Ongoing monitoring procedures established",
        "Actuarial Standards of Practice (ASOPs) Compliance:",
        "  • ASOP 12: Risk Classification - Compliant",
        "  • ASOP 23: Data Quality - Compliant",
        "  • ASOP 41: Actuarial Communications - Compliant",
        "  • ASOP 56: Modeling - Compliant",
        "State Insurance Department Requirements: All 15 states reviewed and approved"
    ])

    # Slide 20: Recommendation
    slide = create_content_slide(prs, "Recommendation & Next Steps")
    tf = slide.placeholders[1].text_frame
    p = tf.paragraphs[0]
    p.text = "RECOMMENDATION:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 119, 180)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Approve the Bodily Injury Frequency Model for production deployment"
    p.font.size = Pt(16)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = ""

    add_bullet_points(tf, [
        "",
        "Model demonstrates superior performance vs. current champion (+22% Gini)",
        "All validation tests passed with excellent results",
        "Regulatory compliance verified across all requirements",
        "Implementation plan is comprehensive and realistic",
        "Monitoring framework will ensure ongoing model effectiveness",
        "",
        "Risk: Low - Model is well-validated and follows industry best practices",
        "Business Impact: High - Improved risk segmentation and rate adequacy"
    ])

    # Save presentation
    output_path = Path("C:/Projects/autodoc-ai/data/examples/bodily_injury_frequency_model.pptx")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"[OK] Saved: {output_path}")

    return output_path


def create_collision_severity_model_ppt():
    """
    Create Collision Severity Model presentation (GLM Gamma + Telematics).

    Based on 2024_collision_model_doc.md:
    - GLM Gamma regression
    - Telematics data included
    - R² 0.52 (+8% vs 2019)
    """
    print("\nCreating Collision Severity Model PPT...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = create_title_slide(
        prs,
        "Collision Severity Model Enhancement",
        "GLM Gamma Regression with Telematics Data\n2024 Model Update\nPresented to Model Risk Committee"
    )

    # Slide 2: Executive Summary
    slide = create_content_slide(prs, "Executive Summary")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Enhanced GLM Gamma model for Collision claim severity prediction",
        "NEW: Incorporates telematics data (hard braking, speeding, night driving)",
        "Model Performance: R² 0.52, MAPE 24.3% on holdout test set",
        "8% improvement vs. 2019 champion model (R² 0.48)",
        "14 statistically significant predictors including 3 telematics variables",
        "Aligns with new 15-territory rating structure (2023 update)",
        "Recommendation: Approve for production deployment Q4 2024"
    ])

    # Slide 3: Business Context
    slide = create_content_slide(prs, "Business Context & Objectives")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Collision coverage represents 35% of total written premium",
        "Parts inflation +18% cumulative since 2019 (supply chain impacts)",
        "Increasing adoption of telematics programs (25% of policyholders)",
        "Opportunity to leverage behavioral driving data for improved pricing",
        "Current model does not reflect ADAS technology impact on severity",
        "Need to align with 2023 territory rerating (10 zones → 15 zones)",
        "Objectives: Improve predictive accuracy, incorporate new data sources"
    ])

    # Slide 4: Data Overview
    slide = create_content_slide(prs, "Data Environment & Sources")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Primary: ClaimsVision system (2020-2023 collision claims)",
        "NEW: Telematics data from DriveWise program (2021-2023)",
        "Sample Size: 450,000 policies, 75,000 claims with severity data",
        "Telematics Coverage: 112,500 policies (25% of total)",
        "Geography: 15 states with updated territory definitions",
        "Severity Range: $500 - $50,000 (excluded total losses > $50K)"
    ])

    # Data summary table
    data = [
        ["Data Source", "Type", "Records", "Time Period", "Quality Score"],
        ["ClaimsVision", "Claims", "75,000", "2020-2023", "99.2%"],
        ["PolicyMaster", "Policy", "450,000", "2020-2023", "98.8%"],
        ["DriveWise", "Telematics", "112,500", "2021-2023", "97.5%"],
        ["VINtelligence", "Vehicle", "450,000", "2020-2023", "99.5%"],
        ["CreditScore DB", "Credit", "405,000", "2020-2023", "96.2%"]
    ]
    add_table(slide, 6, 5, Inches(0.5), Inches(3.8), Inches(9), Inches(2.5), data)

    # Slide 5: Severity Distribution
    slide = create_content_slide(prs, "Claim Severity Descriptive Statistics")
    data = [
        ["Statistic", "Overall", "With Telematics", "Without Telematics"],
        ["Mean Severity", "$3,850", "$3,520", "$4,050"],
        ["Median Severity", "$2,750", "$2,550", "$2,900"],
        ["Std Deviation", "$2,450", "$2,280", "$2,580"],
        ["75th Percentile", "$4,950", "$4,520", "$5,250"],
        ["90th Percentile", "$7,800", "$7,200", "$8,200"],
        ["95th Percentile", "$11,500", "$10,800", "$12,000"],
        ["Max (capped)", "$50,000", "$50,000", "$50,000"]
    ]
    add_table(slide, 8, 4, Inches(1.2), Inches(1.8), Inches(7.6), Inches(4), data)

    tf = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.8)).text_frame
    p = tf.paragraphs[0]
    p.text = "Note: Telematics participants show 13% lower average severity - reflects safer driving behavior"
    p.font.size = Pt(11)
    p.font.italic = True

    # Slide 6: Variable Categories
    slide = create_content_slide(prs, "Candidate Variable Categories")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Vehicle Characteristics (8 variables):",
        "  • Vehicle age, make, model, body type, MSRP, safety rating, ADAS features",
        "Driver Demographics (5 variables):",
        "  • Age, gender, marital status, years licensed, occupation",
        "Geographic (2 variables):",
        "  • Territory (15 zones), garaging ZIP code density",
        "Coverage & Policy (3 variables):",
        "  • Deductible, coverage limits, policy tenure",
        "NEW - Telematics Behavioral (5 variables):",
        "  • Hard braking events, speeding instances, night driving %, mileage, trip count"
    ])

    # Slide 7: Telematics Variables - New Addition
    slide = create_content_slide(prs, "Telematics Variables: Behavioral Risk Factors")
    data = [
        ["Variable", "Definition", "Range", "Mean", "Correlation with Severity"],
        ["Hard Braking Rate", "Events per 100 miles", "0-15", "2.3", "+0.28"],
        ["Speeding Rate", "% miles >10mph over limit", "0%-45%", "8.5%", "+0.31"],
        ["Night Driving %", "% trips 10pm-5am", "0%-60%", "12%", "+0.18"],
        ["Annual Mileage", "Miles driven per year", "1K-40K", "11,200", "+0.12"],
        ["Trip Frequency", "Trips per month", "5-150", "65", "+0.08"]
    ]
    add_table(slide, 6, 5, Inches(0.3), Inches(1.8), Inches(9.4), Inches(3.5), data)

    tf = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.2)).text_frame
    add_bullet_points(tf, [
        "Hard braking and speeding show strongest correlation with claim severity",
        "Night driving associated with higher severity (reduced visibility, fatigue)",
        "Mileage has weak but positive correlation (exposure effect)"
    ])

    # Slide 8: Final Variable Selection
    slide = create_content_slide(prs, "Final Model Variables (14 Predictors)")
    data = [
        ["Variable", "Type", "p-value", "Coefficient", "Impact on Severity"],
        ["Intercept", "-", "< 0.001", "8.125", "Base: $3,350"],
        ["Vehicle Age", "Continuous", "< 0.001", "0.028", "+2.8% per year"],
        ["Vehicle MSRP", "Continuous", "< 0.001", "0.00045", "+0.045% per $100"],
        ["ADAS Level", "Categorical", "< 0.001", "varies", "-8% to -15%"],
        ["Territory", "Categorical", "< 0.001", "varies", "-12% to +18%"],
        ["Deductible", "Categorical", "< 0.001", "varies", "-22% to base"],
        ["Driver Age", "Continuous", "0.003", "-0.005", "-0.5% per year"],
        ["Vehicle Body Type", "Categorical", "< 0.001", "varies", "-5% to +12%"],
        ["Hard Braking Rate", "Continuous", "< 0.001", "0.035", "+3.6% per event"],
        ["Speeding Rate", "Continuous", "< 0.001", "0.008", "+0.8% per % point"],
        ["Night Driving %", "Continuous", "0.002", "0.004", "+0.4% per % point"],
        ["Safety Rating", "Categorical", "0.018", "varies", "-3% to +5%"],
        ["Gender", "Binary", "0.035", "0.042", "Males +4.3%"],
        ["Policy Tenure", "Continuous", "0.041", "-0.012", "-1.2% per year"]
    ]
    add_table(slide, 15, 5, Inches(0.2), Inches(1.3), Inches(9.6), Inches(5.5), data)

    # Slide 9: Single Factor - Vehicle Age
    slide = create_content_slide(prs, "Single Factor Analysis: Vehicle Age")
    data = [
        ["Age Band", "Vehicles", "Claims", "Avg Severity", "Index to Overall", "Significance"],
        ["0-2 years", "85,000", "10,200", "$4,850", "126", "High"],
        ["3-5 years", "120,000", "15,600", "$3,950", "103", "High"],
        ["6-8 years", "105,000", "14,700", "$3,650", "95", "High"],
        ["9-12 years", "90,000", "13,500", "$3,420", "89", "High"],
        ["13+ years", "50,000", "7,500", "$2,950", "77", "Medium"]
    ]
    add_table(slide, 6, 6, Inches(1), Inches(1.8), Inches(8), Inches(3.5), data)

    tf = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1)).text_frame
    add_bullet_points(tf, [
        "Newer vehicles have 26% higher severity (more expensive parts, advanced tech)",
        "Older vehicles show lower severity (aftermarket parts, repair shop flexibility)",
        "Strong monotonic relationship supports continuous variable usage"
    ])

    # Slide 10: Single Factor - ADAS Technology
    slide = create_content_slide(prs, "Single Factor Analysis: ADAS Features (NEW)")
    data = [
        ["ADAS Level", "Description", "Vehicles", "Claims", "Avg Severity", "Index"],
        ["None", "No ADAS features", "180,000", "30,600", "$4,150", "108"],
        ["Basic", "FCW, LDW only", "150,000", "23,250", "$3,820", "99"],
        ["Advanced", "AEB, ACC, BSM", "120,000", "16,800", "$3,250", "84"]
    ]
    add_table(slide, 4, 6, Inches(0.8), Inches(2), Inches(8.4), Inches(2.5), data)

    tf = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(8.4), Inches(1.8)).text_frame
    add_bullet_points(tf, [
        "Advanced ADAS reduces severity by 16% vs. no features",
        "Reflects collision avoidance technology effectiveness",
        "Growing trend: 27% of fleet has ADAS, up from 15% in 2019",
        "Important pricing signal for encouraging safety technology adoption"
    ])

    # Slide 11: Single Factor - Telematics Hard Braking
    slide = create_content_slide(prs, "Single Factor Analysis: Hard Braking Events (Telematics)")
    data = [
        ["Braking Events/100mi", "Drivers", "Claims", "Avg Severity", "Index", "Category"],
        ["0-1 (Low)", "45,000", "5,850", "$3,150", "82", "Safe"],
        ["1.1-2 (Below Avg)", "32,500", "4,550", "$3,520", "91", "Good"],
        ["2.1-3 (Average)", "22,000", "3,300", "$3,850", "100", "Average"],
        ["3.1-5 (Above Avg)", "10,000", "1,700", "$4,420", "115", "Risky"],
        ["5+ (High)", "3,000", "600", "$5,250", "136", "Very Risky"]
    ]
    add_table(slide, 6, 6, Inches(0.6), Inches(1.8), Inches(8.8), Inches(3.5), data)

    tf = slide.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(8.8), Inches(1)).text_frame
    add_bullet_points(tf, [
        "Clear severity gradient: High braking drivers show 36% higher severity",
        "Reflects aggressive driving behavior leading to higher-speed collisions",
        "Strong predictive value supports telematics program expansion"
    ])

    # Slide 12: Model Development
    slide = create_content_slide(prs, "Model Development Approach")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Model Type: GLM with Gamma distribution, log link function",
        "Gamma distribution appropriate for continuous, positive severity values",
        "Software: R glm() function with family=Gamma(link='log')",
        "Data Split: 60% training / 20% validation / 20% test",
        "Missing Data: Telematics missing for 75% of policies (not random)",
        "  • Created separate 'No Telematics' category (not MCAR)",
        "  • Allows model to use telematics when available, neutral otherwise",
        "Outlier Treatment: Capped severity at $50K (excluded total losses)",
        "Validation: Dispersion parameter validated (φ = 1.15, reasonable)"
    ])

    # Slide 13: Champion vs. Challenger
    slide = create_content_slide(prs, "Model Comparison: 2019 Champion vs. 2024 Challenger")
    data = [
        ["Metric", "Champion (2019)", "Challenger (2024)", "Change", "Assessment"],
        ["R-Squared", "0.48", "0.52", "+8.3%", "✓ Better"],
        ["MAPE", "26.8%", "24.3%", "-2.5 pts", "✓ Better"],
        ["RMSE", "$2,850", "$2,620", "-8.1%", "✓ Better"],
        ["Top Quintile R²", "0.41", "0.48", "+17.1%", "✓ Better"],
        ["Deviance", "285,420", "271,850", "-4.8%", "✓ Better"],
        ["AIC", "892,150", "878,920", "-1.5%", "✓ Better"],
        ["# Variables", "11", "14", "+3", "More complex"],
        ["Telematics", "No", "Yes", "NEW", "Innovation"],
        ["ADAS", "No", "Yes", "NEW", "Innovation"],
        ["Territory Zones", "10", "15", "+5", "Updated"]
    ]
    add_table(slide, 11, 5, Inches(0.5), Inches(1.5), Inches(9), Inches(5), data)

    # Slide 14: Model Coefficients
    slide = create_content_slide(prs, "Key Model Coefficients & Interpretations")
    data = [
        ["Variable", "Coefficient", "Exp(Coef)", "Interpretation"],
        ["Intercept", "8.125", "$3,350", "Base severity"],
        ["Vehicle Age (per year)", "0.028", "1.028", "+2.8% per year older"],
        ["ADAS: Advanced vs None", "-0.165", "0.848", "-15.2% lower severity"],
        ["Hard Braking (per event)", "0.035", "1.036", "+3.6% per event/100mi"],
        ["Speeding Rate (per %)", "0.008", "1.008", "+0.8% per percentage point"],
        ["Deductible: $1000 vs $250", "-0.248", "0.780", "-22.0% lower severity"],
        ["Territory: Urban vs Rural", "0.185", "1.203", "+20.3% higher severity"],
        ["Vehicle MSRP ($10K)", "0.045", "1.046", "+4.6% per $10K MSRP"]
    ]
    add_table(slide, 9, 4, Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5), data)

    # Slide 15: Validation Results
    slide = create_content_slide(prs, "Out-of-Sample Validation Performance")
    data = [
        ["Segment", "Claims", "Actual Avg", "Predicted Avg", "A/P Ratio", "MAPE"],
        ["Overall", "15,000", "$3,850", "$3,875", "0.99", "24.3%"],
        ["With Telematics", "3,750", "$3,520", "$3,545", "0.99", "22.8%"],
        ["Without Telematics", "11,250", "$4,050", "$4,070", "1.00", "25.1%"],
        ["ADAS Advanced", "2,400", "$3,250", "$3,280", "0.99", "21.5%"],
        ["ADAS Basic", "4,650", "$3,820", "$3,850", "0.99", "23.8%"],
        ["No ADAS", "7,950", "$4,150", "$4,180", "0.99", "25.7%"],
        ["Deductible $250", "4,500", "$4,520", "$4,550", "0.99", "26.2%"],
        ["Deductible $500", "6,000", "$3,750", "$3,775", "0.99", "24.1%"],
        ["Deductible $1000", "4,500", "$3,250", "$3,270", "0.99", "22.5%"]
    ]
    add_table(slide, 10, 6, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5), data)

    # Slide 16: Residual Analysis
    slide = create_content_slide(prs, "Model Diagnostics & Residual Analysis")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Residual Distribution: Normal Q-Q plot shows good fit (slight heavy tails)",
        "Homoscedasticity: Residuals vs. fitted shows constant variance",
        "No systematic patterns in residuals by key segments",
        "Dispersion Parameter φ = 1.15 (acceptable for Gamma GLM)",
        "Cook's Distance: No influential outliers detected",
        "VIF Analysis: All VIF < 3.5 (multicollinearity not a concern)",
        "Telematics variables show no correlation with non-telematics predictors",
        "Model assumptions validated across all diagnostic tests"
    ])

    # Slide 17: Implementation Plan
    slide = create_content_slide(prs, "Implementation Plan & Timeline")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Q3 2024: Model approval and regulatory filing preparation",
        "Q4 2024: Production deployment",
        "  • October: System integration and testing",
        "  • November: Parallel run with 2019 model",
        "  • December: Production cutover",
        "Telematics Integration:",
        "  • DriveWise API real-time connection",
        "  • Nightly batch updates for behavioral scores",
        "  • Fallback logic when telematics unavailable",
        "Monitoring Dashboard: Live tracking of A/P ratios by segment",
        "Q1 2025: 90-day post-implementation review"
    ])

    # Slide 18: Monitoring Framework
    slide = create_content_slide(prs, "Ongoing Monitoring & Performance Tracking")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Monthly KPIs:",
        "  • Actual vs. Predicted severity by territory, vehicle age, ADAS level",
        "  • Telematics data availability % (target: maintain 25%+)",
        "  • Model drift indicators (PSI < 0.25 threshold)",
        "Quarterly Deep Dive:",
        "  • Segmented A/P analysis (20+ segments)",
        "  • New vs. renewal business comparison",
        "  • Telematics participants vs. non-participants",
        "Annual Refresh:",
        "  • Full model recalibration with latest year data",
        "  • Assess need for new variables or interactions",
        "  • Regulatory validation and audit"
    ])

    # Slide 19: Compliance
    slide = create_content_slide(prs, "Regulatory Compliance & Risk Assessment")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Regulatory Status: APPROVED by all 15 state insurance departments",
        "  • Telematics usage approved in all jurisdictions",
        "  • ADAS variables reviewed for correlation/causation",
        "  • No unfair discrimination concerns identified",
        "ASOP Compliance:",
        "  • ASOP 12 (Risk Classification): Compliant - all variables actuarially justified",
        "  • ASOP 23 (Data Quality): Compliant - telematics quality validated",
        "  • ASOP 41 (Communications): Compliant - documentation complete",
        "  • ASOP 56 (Modeling): Compliant - appropriate model selection",
        "Model Risk Rating: MODERATE (due to telematics data dependency)",
        "Mitigation: Robust fallback logic when telematics unavailable"
    ])

    # Slide 20: Recommendation
    slide = create_content_slide(prs, "Recommendation & Business Impact")
    tf = slide.placeholders[1].text_frame
    p = tf.paragraphs[0]
    p.text = "RECOMMENDATION:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 119, 180)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Approve Collision Severity Model for Q4 2024 production deployment"
    p.font.size = Pt(16)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = ""

    add_bullet_points(tf, [
        "",
        "Model Performance: 8% improvement in R² vs. champion (+17% in top quintile)",
        "Innovation: First model to leverage telematics behavioral data",
        "Risk Segmentation: ADAS and telematics provide material differentiation",
        "Validation: Excellent A/P ratios across all segments (0.99-1.00)",
        "Regulatory: Full compliance, approved by all 15 states",
        "",
        "Expected Impact:",
        "  • Improved pricing accuracy → better loss ratios",
        "  • Telematics adoption incentive → safer driving behaviors",
        "  • Competitive advantage in ADAS-equipped vehicle segment"
    ])

    # Save presentation
    output_path = Path("C:/Projects/autodoc-ai/data/examples/collision_severity_model.pptx")
    prs.save(str(output_path))
    print(f"[OK] Saved: {output_path}")

    return output_path


def create_comprehensive_coverage_model_ppt():
    """
    Create Comprehensive Coverage Model presentation (XGBoost).

    Based on 2023_comprehensive_model_doc.md:
    - XGBoost with 42 features
    - SHAP explainability
    - Frequency R² 0.58, Severity R² 0.44
    """
    print("\nCreating Comprehensive Coverage Model PPT...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = create_title_slide(
        prs,
        "Comprehensive Coverage Model",
        "XGBoost Machine Learning Model with SHAP Explainability\n2024 Model Development\nPresented to Model Risk Committee"
    )

    # Slide 2: Executive Summary
    slide = create_content_slide(prs, "Executive Summary")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Developed XGBoost gradient boosting model for Comprehensive coverage",
        "Dual prediction: Frequency (R² 0.58) and Severity (R² 0.44)",
        "42 predictor variables with engineered features and interactions",
        "38% improvement in frequency prediction vs. GLM benchmark (R² 0.42)",
        "42% improvement in severity prediction vs. GLM benchmark (R² 0.31)",
        "SHAP values provide model explainability for regulatory acceptance",
        "Successfully addresses 'black box' concerns through transparency framework",
        "Recommendation: Approve for production deployment Q1 2025"
    ])

    # Slide 3: Business Context
    slide = create_content_slide(prs, "Business Context & Strategic Rationale")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Comprehensive coverage: $2.1B annual premium, 28% of total book",
        "Current GLM models underperforming in competitive markets",
        "Theft trends changing rapidly (catalytic converter, keyless entry exploits)",
        "Weather patterns becoming less predictable (climate change impact)",
        "Opportunity to leverage machine learning for complex interaction effects",
        "Regulatory environment increasingly accepting of ML with explainability",
        "Objectives:",
        "  • Improve risk segmentation by 30-40% vs. current models",
        "  • Maintain full transparency and explainability",
        "  • Address emerging theft and weather-related risks"
    ])

    # Slide 4: Data Environment
    slide = create_content_slide(prs, "Data Sources & Coverage")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Primary Claims Data: ClaimsVision (2019-2023)",
        "Policy Data: PolicyMaster (2019-2023)",
        "External Data Sources (NEW):",
        "  • NICB Theft Database (vehicle theft risk scores by make/model/ZIP)",
        "  • NOAA Weather (historical hail, wind, flood events by ZIP)",
        "  • Garaging location: Secure vs. street parking (from underwriting)",
        "Sample Size: 600,000 policies, 42,000 comprehensive claims",
        "Time Period: 5 years (2019-2023)",
        "Data Quality: 97.8% complete after quality checks"
    ])

    # Data summary
    data = [
        ["Metric", "Training", "Validation", "Test"],
        ["Policies", "360,000", "120,000", "120,000"],
        ["Frequency Claims", "21,600", "7,200", "7,200"],
        ["Severity Claims", "18,000", "6,000", "6,000"],
        ["Avg Frequency", "6.0%", "6.0%", "6.0%"],
        ["Avg Severity", "$2,850", "$2,820", "$2,875"],
        ["Time Period", "2019-2021", "2022", "2023"]
    ]
    add_table(slide, 7, 4, Inches(0.8), Inches(4), Inches(8.4), Inches(2.8), data)

    # Slide 5: Why Machine Learning?
    slide = create_content_slide(prs, "Rationale for XGBoost vs. Traditional GLM")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Complex Interaction Effects:",
        "  • Weather + vehicle age + garaging type interactions",
        "  • Theft risk varies non-linearly by ZIP + vehicle make + security features",
        "  • Traditional GLM limited to 2-way interactions (model becomes unwieldy)",
        "Non-Linear Relationships:",
        "  • Hail damage threshold effects (size-dependent damage)",
        "  • Theft risk curves (not linear with vehicle value)",
        "Automatic Feature Engineering:",
        "  • XGBoost discovers optimal splits and interactions",
        "  • Reduces manual feature engineering burden",
        "Superior Predictive Performance:",
        "  • Demonstrated 38-42% improvement in pilot testing",
        "Explainability via SHAP:",
        "  • Addresses regulatory 'black box' concerns",
        "  • Individual prediction explanations available"
    ])

    # Slide 6: Variable Categories (42 variables)
    slide = create_content_slide(prs, "Variable Categories: 42 Total Predictors")
    data = [
        ["Category", "# Variables", "Examples", "Data Source"],
        ["Vehicle", "12", "Make, model, year, value, theft rating, security", "VINtelligence"],
        ["Driver/Policy", "8", "Age, gender, credit, tenure, deductible", "PolicyMaster"],
        ["Geographic", "6", "Territory, ZIP, urban/rural, garage type", "PolicyMaster"],
        ["Weather Risk", "8", "Hail freq, wind speed, flood zone, precip", "NOAA"],
        ["Theft Risk", "5", "NICB score, local crime, vehicle theft rank", "NICB + FBI"],
        ["Engineered", "3", "Value*Theft, Weather*Age, Garage*Territory", "Created"]
    ]
    add_table(slide, 7, 4, Inches(0.8), Inches(2), Inches(8.4), Inches(3.5), data)

    tf = slide.shapes.add_textbox(Inches(0.8), Inches(5.7), Inches(8.4), Inches(1)).text_frame
    add_bullet_points(tf, [
        "Mix of traditional actuarial variables + external data",
        "Engineered features capture domain knowledge",
        "All variables validated for actuarial soundness and non-discrimination"
    ])

    # Slide 7: Top 15 Variables by SHAP Importance
    slide = create_content_slide(prs, "Top 15 Variables by SHAP Feature Importance")
    data = [
        ["Rank", "Variable", "SHAP Mean |Value|", "Type", "Predictive Focus"],
        ["1", "Vehicle Value", "0.185", "Continuous", "Frequency + Severity"],
        ["2", "NICB Theft Score", "0.142", "Continuous", "Frequency (Theft)"],
        ["3", "Territory", "0.128", "Categorical", "Frequency + Severity"],
        ["4", "Hail Frequency (ZIP)", "0.115", "Continuous", "Frequency + Severity"],
        ["5", "Vehicle Age", "0.098", "Continuous", "Frequency + Severity"],
        ["6", "Deductible", "0.087", "Categorical", "Severity"],
        ["7", "Garaging Type", "0.076", "Categorical", "Frequency (Theft)"],
        ["8", "Vehicle Make (Luxury)", "0.068", "Binary", "Severity"],
        ["9", "Credit Score", "0.062", "Continuous", "Frequency"],
        ["10", "Flood Zone", "0.055", "Binary", "Frequency"],
        ["11", "Driver Age", "0.048", "Continuous", "Frequency"],
        ["12", "Vehicle Security Features", "0.045", "Categorical", "Frequency (Theft)"],
        ["13", "Wind Speed (Historical)", "0.041", "Continuous", "Frequency + Severity"],
        ["14", "Urban Density", "0.038", "Continuous", "Frequency (Theft)"],
        ["15", "Policy Tenure", "0.032", "Continuous", "Frequency"]
    ]
    add_table(slide, 16, 5, Inches(0.2), Inches(1.3), Inches(9.6), Inches(5.7), data)

    # Slide 8: Single Factor - Vehicle Value
    slide = create_content_slide(prs, "Single Factor Analysis: Vehicle Value")
    data = [
        ["Value Band", "Vehicles", "Freq Claims", "Claim Freq", "Avg Severity", "Pure Premium"],
        ["<$10K", "120,000", "5,400", "4.5%", "$1,850", "$83"],
        ["$10K-$20K", "180,000", "9,900", "5.5%", "$2,350", "$129"],
        ["$20K-$30K", "150,000", "9,750", "6.5%", "$2,950", "$192"],
        ["$30K-$50K", "100,000", "7,500", "7.5%", "$3,850", "$289"],
        ["$50K+", "50,000", "4,500", "9.0%", "$5,250", "$473"]
    ]
    add_table(slide, 6, 6, Inches(0.8), Inches(2), Inches(8.4), Inches(3.5), data)

    tf = slide.shapes.add_textbox(Inches(0.8), Inches(5.7), Inches(8.4), Inches(1)).text_frame
    add_bullet_points(tf, [
        "Strong relationship: Higher value → higher frequency AND severity",
        "Pure premium increases 470% from lowest to highest band",
        "Reflects both theft targeting (frequency) and repair costs (severity)"
    ])

    # Slide 9: Single Factor - NICB Theft Score
    slide = create_content_slide(prs, "Single Factor Analysis: NICB Theft Risk Score")
    data = [
        ["Theft Score", "Description", "Vehicles", "Freq Claims", "Frequency", "Index"],
        ["1-2", "Very Low Risk", "150,000", "6,000", "4.0%", "67"],
        ["3-4", "Low Risk", "180,000", "9,000", "5.0%", "83"],
        ["5-6", "Medium Risk", "150,000", "9,750", "6.5%", "108"],
        ["7-8", "High Risk", "90,000", "7,650", "8.5%", "142"],
        ["9-10", "Very High Risk", "30,000", "3,600", "12.0%", "200"]
    ]
    add_table(slide, 6, 6, Inches(0.8), Inches(2), Inches(8.4), Inches(3.5), data)

    tf = slide.shapes.add_textbox(Inches(0.8), Inches(5.7), Inches(8.4), Inches(1.2)).text_frame
    add_bullet_points(tf, [
        "NICB score is 2nd most important variable (SHAP importance)",
        "3x frequency difference between very low and very high risk",
        "Correlates with specific make/model theft targets (e.g., Honda Accord, Kia/Hyundai)"
    ])

    # Slide 10: Single Factor - Hail Risk
    slide = create_content_slide(prs, "Single Factor Analysis: Hail Frequency by ZIP")
    data = [
        ["Hail Events/Year", "ZIPs", "Vehicles", "Freq Claims", "Frequency", "Avg Severity"],
        ["0", "Low Risk", "200,000", "10,000", "5.0%", "$2,450"],
        ["0.1-0.5", "Low-Medium", "180,000", "10,800", "6.0%", "$2,750"],
        ["0.6-1.5", "Medium", "120,000", "8,400", "7.0%", "$3,150"],
        ["1.6-3.0", "High", "80,000", "6,400", "8.0%", "$3,550"],
        ["3.0+", "Very High", "20,000", "1,800", "9.0%", "$4,250"]
    ]
    add_table(slide, 6, 6, Inches(0.8), Inches(2), Inches(8.4), Inches(3.5), data)

    tf = slide.shapes.add_textbox(Inches(0.8), Inches(5.7), Inches(8.4), Inches(1)).text_frame
    add_bullet_points(tf, [
        "Hail is 4th most important predictor (SHAP importance)",
        "Impacts both frequency (80% increase) and severity (73% increase)",
        "Climate change increasing hail frequency and severity in certain regions"
    ])

    # Slide 11: XGBoost Model Architecture
    slide = create_content_slide(prs, "XGBoost Model Architecture & Hyperparameters")

    data = [
        ["Parameter", "Frequency Model", "Severity Model", "Selection Method"],
        ["Objective", "binary:logistic", "reg:squarederror", "Problem type"],
        ["# Trees", "350", "280", "Early stopping CV"],
        ["Max Depth", "6", "5", "Grid search CV"],
        ["Learning Rate", "0.05", "0.08", "Grid search CV"],
        ["Min Child Weight", "3", "5", "Grid search CV"],
        ["Subsample", "0.8", "0.75", "Grid search CV"],
        ["Colsample (tree)", "0.8", "0.85", "Grid search CV"],
        ["Gamma", "0.1", "0.15", "Grid search CV"],
        ["Reg Alpha (L1)", "0.05", "0.02", "Grid search CV"],
        ["Reg Lambda (L2)", "1.0", "1.5", "Grid search CV"]
    ]
    add_table(slide, 11, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(5), data)

    tf = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9), Inches(0.5)).text_frame
    p = tf.paragraphs[0]
    p.text = "Hyperparameters selected via 5-fold cross-validation with early stopping to prevent overfitting"
    p.font.size = Pt(10)
    p.font.italic = True

    # Slide 12: Frequency Model - Champion vs Challenger
    slide = create_content_slide(prs, "Frequency Model: XGBoost vs. GLM Benchmark")
    data = [
        ["Metric", "GLM Champion", "XGBoost Challenger", "Improvement", "Assessment"],
        ["R-Squared", "0.42", "0.58", "+38.1%", "✓ Significant"],
        ["AUC", "0.68", "0.76", "+11.8%", "✓ Excellent"],
        ["Gini", "0.36", "0.52", "+44.4%", "✓ Significant"],
        ["Top Decile Lift", "2.1x", "2.8x", "+33.3%", "✓ Excellent"],
        ["Accuracy", "94.2%", "94.8%", "+0.6 pts", "✓ Better"],
        ["Precision (Top 20%)", "12.5%", "15.8%", "+26.4%", "✓ Better"],
        ["Log Loss", "0.285", "0.238", "-16.5%", "✓ Better"]
    ]
    add_table(slide, 8, 5, Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5), data)

    tf = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(8.4), Inches(0.7)).text_frame
    p = tf.paragraphs[0]
    p.text = "XGBoost shows substantial improvement across all metrics, particularly in risk discrimination (Gini +44%)"
    p.font.size = Pt(11)
    p.font.bold = True

    # Slide 13: Severity Model - Champion vs Challenger
    slide = create_content_slide(prs, "Severity Model: XGBoost vs. GLM Benchmark")
    data = [
        ["Metric", "GLM Champion", "XGBoost Challenger", "Improvement", "Assessment"],
        ["R-Squared", "0.31", "0.44", "+41.9%", "✓ Significant"],
        ["RMSE", "$1,850", "$1,620", "-12.4%", "✓ Better"],
        ["MAE", "$1,280", "$1,095", "-14.5%", "✓ Better"],
        ["MAPE", "32.5%", "28.2%", "-4.3 pts", "✓ Better"],
        ["Top Quintile R²", "0.28", "0.41", "+46.4%", "✓ Excellent"],
        ["Median APE", "24.8%", "21.2%", "-3.6 pts", "✓ Better"]
    ]
    add_table(slide, 7, 5, Inches(0.8), Inches(1.8), Inches(8.4), Inches(3.8), data)

    tf = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(8.4), Inches(1)).text_frame
    add_bullet_points(tf, [
        "42% improvement in R² demonstrates superior severity prediction",
        "Particularly strong in high-value vehicles and hail-prone areas",
        "Captures non-linear relationships (e.g., hail size thresholds)"
    ])

    # Slide 14: SHAP Explainability
    slide = create_content_slide(prs, "SHAP Values: Model Explainability Framework")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "SHAP (SHapley Additive exPlanations) provides global and local interpretability",
        "Addresses regulatory 'black box' concern with ML models",
        "Individual Prediction Explanation:",
        "  • Every prediction shows contribution of each variable",
        "  • Allows review of outlier predictions",
        "  • Supports underwriting exceptions and pricing discussions",
        "Global Feature Importance:",
        "  • Identifies most influential variables across all predictions",
        "  • Validates actuarial intuition (vehicle value, theft, weather most important)",
        "Regulatory Acceptance:",
        "  • Approved by 12 of 15 state regulators (3 pending)",
        "  • SHAP documentation included in all rate filings"
    ])

    # Slide 15: SHAP Example - High Risk Prediction
    slide = create_content_slide(prs, "SHAP Explanation Example: High Frequency Prediction")

    tf = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.8)).text_frame
    add_bullet_points(tf, [
        "Example Policy: 2021 Honda Accord, $28K value, Chicago ZIP (high theft), street parking, age 25 driver",
        "Predicted Frequency: 11.5% (base rate: 6.0%, prediction +5.5 percentage points)"
    ])

    data = [
        ["Variable", "Value", "SHAP Contribution", "Direction", "Explanation"],
        ["Base Rate", "-", "6.0%", "Base", "Average frequency"],
        ["NICB Theft Score", "9/10", "+2.8%", "Increase", "Honda Accord high theft risk"],
        ["Territory (Chicago)", "Zone 2", "+1.5%", "Increase", "Urban high-density area"],
        ["Garaging Type", "Street", "+0.9%", "Increase", "No secure garage"],
        ["Vehicle Value", "$28K", "+0.6%", "Increase", "Above average value"],
        ["Driver Age", "25", "+0.3%", "Increase", "Young driver factor"],
        ["Vehicle Age", "3 years", "-0.2%", "Decrease", "Not brand new (less theft)"],
        ["Credit Score", "720", "-0.4%", "Decrease", "Good credit mitigates"],
        ["Total Predicted", "-", "11.5%", "Final", "Sum of contributions"]
    ]
    add_table(slide, 9, 5, Inches(0.3), Inches(2.5), Inches(9.4), Inches(4), data)

    # Slide 16: Validation Results
    slide = create_content_slide(prs, "Out-of-Sample Validation: Test Set Performance")
    data = [
        ["Segment", "Policies", "Actual Freq", "Pred Freq", "A/P", "RMSE Severity"],
        ["Overall", "120,000", "6.0%", "6.05%", "0.99", "$1,620"],
        ["High Theft Risk", "24,000", "10.2%", "10.3%", "0.99", "$2,150"],
        ["Low Theft Risk", "36,000", "3.8%", "3.75%", "1.01", "$1,420"],
        ["High Hail Risk", "18,000", "8.5%", "8.6%", "0.99", "$3,280"],
        ["Low Hail Risk", "42,000", "4.8%", "4.75%", "1.01", "$1,380"],
        ["Luxury Vehicles", "12,000", "7.5%", "7.6%", "0.99", "$4,850"],
        ["Economy Vehicles", "54,000", "5.2%", "5.15%", "1.01", "$1,250"],
        ["Urban Territory", "48,000", "7.2%", "7.3%", "0.99", "$2,050"],
        ["Rural Territory", "30,000", "4.5%", "4.45%", "1.01", "$1,580"]
    ]
    add_table(slide, 10, 6, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5), data)

    # Slide 17: Overfitting Analysis
    slide = create_content_slide(prs, "Overfitting Prevention & Cross-Validation")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Early Stopping: Training stopped when validation error stops improving",
        "Cross-Validation: 5-fold CV during hyperparameter tuning",
        "Regularization: L1 (alpha) and L2 (lambda) penalties applied",
        "Performance Comparison:",
        "  • Training R²: 0.62 (frequency), 0.48 (severity)",
        "  • Validation R²: 0.59 (frequency), 0.45 (severity)",
        "  • Test R²: 0.58 (frequency), 0.44 (severity)",
        "  • Minimal degradation from train → validation → test",
        "Feature Importance Stability:",
        "  • Top 10 variables consistent across all CV folds",
        "  • SHAP values stable across bootstrap samples",
        "Conclusion: No evidence of overfitting, model generalizes well"
    ])

    # Slide 18: Implementation Plan
    slide = create_content_slide(prs, "Implementation Roadmap")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Q4 2024: Final Model Approval & Regulatory Filings",
        "  • Model Risk Committee approval (this presentation)",
        "  • Submit rate filings to 15 states (3 already approved)",
        "Q1 2025: System Implementation",
        "  • Deploy XGBoost scoring engine to production environment",
        "  • Real-time API for quote/bind process (<500ms response time)",
        "  • Batch processing for policy renewals",
        "  • SHAP explanation generation for all policies",
        "Q1 2025: Monitoring & Validation",
        "  • Shadow mode: Run alongside GLM for 30 days",
        "  • A/B test: 10% of quotes use XGBoost pricing",
        "Q2 2025: Full Production Rollout",
        "  • 100% of new business and renewals",
        "  • Monthly performance reviews for first 6 months"
    ])

    # Slide 19: Risk Mitigation & Governance
    slide = create_content_slide(prs, "Model Risk Mitigation & Ongoing Governance")
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, [
        "Model Risk Rating: MODERATE-HIGH (due to ML complexity)",
        "Mitigation Strategies:",
        "  • SHAP explanations for all predictions (transparency)",
        "  • Champion-challenger A/B testing (GLM remains active backup)",
        "  • Prediction bounds: Flag predictions >3 std dev for manual review",
        "  • Monthly model monitoring: PSI, feature drift, performance metrics",
        "  • Independent validation: External actuarial firm review completed",
        "Governance Framework:",
        "  • Quarterly Model Risk Committee review",
        "  • Annual full revalidation and recalibration",
        "  • Feature importance monitoring (alert if top 5 variables change >20%)",
        "  • Explainability audit: Random sample of 100 policies/month reviewed",
        "Regulatory Engagement:",
        "  • Ongoing dialogue with 3 pending state regulators",
        "  • Annual presentation to regulators with model performance"
    ])

    # Slide 20: Recommendation
    slide = create_content_slide(prs, "Recommendation & Strategic Value")
    tf = slide.placeholders[1].text_frame
    p = tf.paragraphs[0]
    p.text = "RECOMMENDATION:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 119, 180)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Approve XGBoost Comprehensive Coverage Model for Q1 2025 deployment"
    p.font.size = Pt(16)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = ""

    add_bullet_points(tf, [
        "",
        "Exceptional Performance: 38-42% improvement vs. GLM benchmark",
        "Regulatory Compliance: SHAP explainability addresses transparency concerns",
        "Risk Management: Comprehensive monitoring and governance framework",
        "Validation: Excellent test set performance across all segments",
        "Competitive Advantage: First-mover in ML for comprehensive coverage",
        "",
        "Expected Business Impact:",
        "  • Improved loss ratios: 3-5 points (better risk selection)",
        "  • Premium growth: Better pricing in high-theft, high-weather-risk markets",
        "  • Market differentiation: ML sophistication attracts tech-savvy customers",
        "",
        "Strategic Value: Establishes ML capabilities for future product lines"
    ])

    # Save presentation
    output_path = Path("C:/Projects/autodoc-ai/data/examples/comprehensive_coverage_model.pptx")
    prs.save(str(output_path))
    print(f"[OK] Saved: {output_path}")

    return output_path


def main():
    """Create all 3 example presentations."""
    print("=" * 60)
    print("Creating Example PowerPoint Presentations for AutoDoc AI")
    print("=" * 60)

    # Create all presentations
    ppt1 = create_frequency_model_ppt()
    ppt2 = create_collision_severity_model_ppt()
    ppt3 = create_comprehensive_coverage_model_ppt()

    print("\n" + "=" * 60)
    print("[OK] All presentations created successfully!")
    print("=" * 60)
    print(f"\n1. {ppt1}")
    print(f"2. {ppt2}")
    print(f"3. {ppt3}")
    print("\nThese presentations can now be used as inputs to AutoDoc AI.")
    print("Upload them through the Streamlit app to test the system!")


if __name__ == "__main__":
    main()
